from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "docs" / "cursor-training-deck.pptx"
FALLBACK_OUTPUT = ROOT / "docs" / "cursor-training-deck-updated.pptx"


TITLE_COLOR = RGBColor(19, 33, 68)
TEXT_COLOR = RGBColor(39, 48, 66)
ACCENT = RGBColor(49, 102, 245)
LIGHT = RGBColor(240, 244, 252)


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 251, 255)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.2), Inches(0.8))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(11.2), Inches(0.5))
        sub_frame = sub_box.text_frame
        sub_p = sub_frame.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.name = "Aptos"
        sub_run.font.size = Pt(12)
        sub_run.font.color.rgb = ACCENT


def add_bullets(slide, items: list[str], left: float = 0.95, top: float = 1.65, width: float = 11.0, height: float = 4.8) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True

    first = True
    for item in items:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(21)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)


def add_prompt_box(slide, prompt: str, top: float = 1.8, height: float = 4.6) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(top),
        Inches(11.4),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1.25)

    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = prompt
    run.font.name = "Consolas"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(28, 34, 52)


def add_steps_table(
    slide,
    rows: list[tuple[str, str, str]],
    top: float = 1.65,
    height: float = 4.95,
) -> None:
    shape = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.6), Inches(top), Inches(12.0), Inches(height))
    table = shape.table

    table.columns[0].width = Inches(1.15)
    table.columns[1].width = Inches(2.45)
    table.columns[2].width = Inches(8.4)

    headers = ["Step", "Goal", "Prompt"]
    for idx, label in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

    for row_idx, (step, goal, prompt) in enumerate(rows, start=1):
        values = [step, goal, prompt]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 else LIGHT
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Aptos" if col_idx < 2 else "Consolas"
                    run.font.size = Pt(12.5 if col_idx < 2 else 10.5)
                    run.font.color.rgb = TEXT_COLOR


def add_feature_table(
    slide,
    rows: list[tuple[str, str, str, str]],
    top: float = 1.65,
    height: float = 4.8,
) -> None:
    shape = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.35), Inches(top), Inches(12.65), Inches(height))
    table = shape.table

    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(3.15)
    table.columns[2].width = Inches(5.1)
    table.columns[3].width = Inches(3.2)

    headers = ["Feature", "Best Practice", "Prompt Or Instructions", "How To Use It In Cursor"]
    for idx, label in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(12.5)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

    for row_idx, (feature, practice, prompt_or_instructions, usage) in enumerate(rows, start=1):
        values = [feature, practice, prompt_or_instructions, usage]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 else LIGHT
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Consolas" if col_idx == 2 else "Aptos"
                    run.font.size = Pt(9.2 if col_idx == 2 else 9.8)
                    run.font.color.rgb = TEXT_COLOR


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.65), Inches(11.2), Inches(0.3))
    frame = box.text_frame
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(100, 110, 130)


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(0.55),
        Inches(5.2),
        Inches(0.55),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT
    banner.line.color.rgb = ACCENT
    banner.text_frame.text = "Cursor Training Deck"
    banner.text_frame.paragraphs[0].font.name = "Aptos"
    banner.text_frame.paragraphs[0].font.size = Pt(16)
    banner.text_frame.paragraphs[0].font.bold = True
    banner.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    banner.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(8.2), Inches(1.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Safe Onramp to Cursor"
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR
    p2 = tf.add_paragraph()
    p2.text = "A 60-minute live walkthrough for understanding code, planning a change, implementing it, and reviewing it safely."
    p2.font.name = "Aptos"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_COLOR
    p2.space_before = Pt(10)

    add_footer(slide, "Event Stream Processor Training Sample")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Session Goal")
    add_bullets(
        slide,
        [
            "Show a practical, low-risk way to use Cursor on a real-looking codebase.",
            "Use AI to improve understanding first, then make one small, constrained feature change.",
            "Verify and review the result before trusting it.",
            "Reinforce that the goal is faster iteration without lowering engineering standards.",
        ],
    )
    add_footer(slide, "1")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Codebase At A Glance", "What to point at in the repo")
    add_bullets(
        slide,
        [
            "sample-data/aws-cloudtrail-login.json: raw AWS-style input event",
            "mappings/aws-to-ecs.yaml: source field to normalized field mapping",
            "backend/src/mapper.py and backend/src/validator.py: transform and validate the payload",
            "backend/src/server.py: serves /api/demo for the frontend",
            "frontend/src/App.jsx: visualizes raw data, mapped output, and validation status",
            "tests/test_mapping.py: current regression coverage",
        ],
    )
    add_footer(slide, "2")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Opening Script")
    add_bullets(
        slide,
        [
            "This sample app is a tiny event-processing pipeline.",
            "The backend loads a raw AWS-style security event, maps it into a normalized schema, validates it, and exposes it through a local API.",
            "The frontend calls that API and shows the raw event, mapped output, and validation issues.",
            "We are using it as a safe sandbox for learning Cursor workflows, not as a production architecture example.",
        ],
        top=1.8,
    )
    add_footer(slide, "3")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Workflow Table", "Steps 1-4 with copy/paste prompts")
    add_steps_table(
        slide,
        [
            (
                "1",
                "Understand the repo",
                "Explain this codebase like I am a new engineer joining the team.\n"
                "Focus on the end-to-end flow from raw event to frontend display.\n"
                "List the 5-8 most important files and what each one does.\n"
                "Keep it concise and practical.",
            ),
            (
                "2",
                "Create a project rule",
                "Create one concise Cursor rule for this repo.\n"
                "The rule should always apply.\n"
                "Focus on non-obvious guidance only: edge cases, recurring mistakes, and places where this repo differs from normal standards.\n"
                "Do not include obvious generic advice like 'write clean code' or 'run tests.'\n"
                "Keep it short enough to help every prompt without bogging the agent down, include one concrete repo example, and explain how the rule will influence future prompts.",
            ),
            (
                "3",
                "Plan the feature",
                "Plan a minimal feature for this sample app.\n"
                "Add `cloud.account.id` to the normalized event by mapping it from `userIdentity.accountId`, "
                "make it required in the schema, update the test coverage, and add a small frontend summary "
                "section that shows `user.name`, `source.ip`, and `cloud.account.id`.\n"
                "Constraints: keep the change small, touch only necessary files, avoid unrelated architecture "
                "or styling changes, tell me which files you expect to edit and why, and include how you will verify the change.",
            ),
            (
                "4",
                "Implement the feature",
                "Implement the planned feature.\n"
                "Add `cloud.account.id` to the normalized event by mapping it from `userIdentity.accountId`, "
                "make it a required schema field, update the test, and add a small frontend summary section "
                "that shows `user.name`, `source.ip`, and `cloud.account.id`.\n"
                "Constraints: keep the implementation simple, only edit the minimum necessary files, do not "
                "refactor unrelated code, and run the relevant test/build checks afterward.",
            ),
        ],
        top=1.55,
        height=5.2,
    )
    add_footer(slide, "4")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Workflow Table", "Steps 5-8 with copy/paste prompts")
    add_steps_table(
        slide,
        [
            (
                "5",
                "Verify the change",
                "Show me how you verified this change.\n"
                "Run the relevant checks and summarize which tests or builds you ran, whether they passed or failed, and any risks or gaps that still remain.\n"
                "Keep it brief.",
            ),
            (
                "6",
                "Inventory changed files",
                "List every file you changed for this feature.\n"
                "For each file, give me one sentence on why it was touched.\n"
                "Also tell me whether any files were intentionally left unchanged.",
            ),
            (
                "7",
                "Ask for self-review",
                "Review this change like a skeptical teammate.\n"
                "Look for bugs, edge cases, unclear naming, accidental scope creep, and missing tests.\n"
                "List findings first, ordered by severity.\n"
                "Then give a short summary of whether this is ready for review.",
            ),
            (
                "8",
                "Close with summary",
                "Summarize what changed in plain English for the team.\n"
                "Explain the feature, the files touched, how it was verified, and why the change is reasonably safe.\n"
                "Keep it short enough that I can read it aloud.",
            ),
        ],
        top=1.55,
        height=5.15,
    )
    add_footer(slide, "5")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Rules, Skills, Subagents", "Configured from the same Cursor panel")
    add_feature_table(
        slide,
        [
            (
                "Rule",
                "Use a rule for a repeated workflow dependency that is easy to miss. Here, mapping changes often require synchronized updates across mapping config, schema, tests, and the demo UI.",
                "Step by step:\n1. Open Rules.\n2. Click `+ New`.\n3. Choose `Project Rule`.\n4. Name it `training-sample-guidance`.\n5. Paste:\n---\ndescription: Keep mapping changes synchronized across config, schema, tests, and demo UI\nalwaysApply: true\n---\n# Mapping Change Workflow\nWhen adding, removing, renaming, or remapping a normalized field:\n- Update `mappings/aws-to-ecs.yaml`.\n- Check whether the field also needs to change in `schemas/ecs-required-fields.json`.\n- Update `tests/test_mapping.py`.\n- If the frontend shows the field or its validation output, update `frontend/src/App.jsx`.\n- Do not stop after changing only the mapping file if schema, tests, or UI would now be inconsistent.\n## Verification\n- Run `python -m pytest ../tests` from `backend/`.\n- If frontend behavior changed, run `npm run build` from `frontend/`.",
                "Create it from the Rules UI as a Project Rule. After that, Cursor applies it automatically on future prompts because `alwaysApply: true` is set.",
            ),
            (
                "Skill",
                "Use a skill for a repeated implementation workflow. Here, the reusable workflow is making mapping changes completely: trace the field, update the related files, run checks, and summarize risks.",
                "Step by step:\n1. Open Skills.\n2. Click `+ New`.\n3. Cursor opens a new agent with `/create-skill Help me create this skill for Cursor:`.\n4. Paste this after that starter text:\nCreate a project skill named `event-mapping-change`.\nUse it when adding, removing, renaming, or reviewing field mappings.\nTell Cursor to:\n- trace the field through sample data, mapping config, transform logic, schema, tests, and frontend\n- prefer small synchronized changes\n- avoid broad refactors unless asked\n- run `python -m pytest ../tests` from `backend/`\n- run `npm run build` from `frontend/` if UI changed\n- report changed files and remaining risks",
                "Create it from the Skills UI. After it exists, Cursor can call it when relevant, or the user can invoke it with `/event-mapping-change`.",
            ),
            (
                "Subagent",
                "Use a subagent when the task is exploratory and benefits from parallel context gathering. Here, the cleanest example is tracing the backend path and frontend display path at the same time before planning a change.",
                "Copy/paste in chat:\nUse two explore subagents in parallel.\nOne should trace the backend flow from `sample-data/aws-cloudtrail-login.json` through `mappings/aws-to-ecs.yaml`, `backend/src/mapper.py`, `backend/src/validator.py`, `backend/src/server.py`, and `/api/demo`.\nThe other should trace how `frontend/src/App.jsx` fetches and displays that data.\nReturn the key files and a short summary of each flow, then synthesize the results for me.",
                "Ask for this directly in chat when you want parallel exploration. Cursor may also decide to use subagents on larger research tasks.",
            ),
        ],
        top=1.45,
        height=4.95,
    )
    add_prompt_box(
        slide,
        "Rules shape behavior. Skills teach workflows. Subagents split broader work into focused helpers.",
        top=6.02,
        height=0.38,
    )
    add_footer(slide, "6")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Verification Commands", "Run checks in Cursor's integrated terminal or ask Cursor to run them")
    add_bullets(
        slide,
        [
            "Use commands for repeatable verification steps after mapping, schema, or UI changes.",
            "Backend check: run `python -m pytest ../tests` from `backend/`.",
            "Frontend check: if UI changed, run `npm run build` from `frontend/`.",
            "Then ask Cursor to summarize what passed, what failed, and any remaining risks.",
        ],
        top=1.75,
        height=2.4,
    )
    add_prompt_box(
        slide,
        "Summarize the verification results from the commands I just ran.\n"
        "Tell me what passed, what failed, and any remaining risks or gaps for review.",
        top=4.15,
        height=1.35,
    )
    add_footer(slide, "7")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Presenter Notes")
    add_bullets(
        slide,
        [
            "Start with understanding, not code generation.",
            "Add one lightweight rule before the feature work so the agent inherits repo-specific judgment.",
            "Use Plan mode when you want guardrails before edits happen.",
            "Keep the live feature small so the audience can follow every touched file.",
            "Model the habit: trust, but verify.",
        ],
        top=1.75,
        height=2.2,
    )
    add_prompt_box(
        slide,
        "Suggested line to say:\n"
        "We are not asking Cursor to do everything. We are giving it a small task with constraints, "
        "then checking what changed before we trust it.",
        top=3.45,
        height=1.6,
    )
    add_footer(slide, "8")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Checklist To Use Live")
    add_bullets(
        slide,
        [
            "1. Explain the repo",
            "2. Create a lightweight project rule",
            "3. Plan the feature",
            "4. Implement the feature",
            "5. Verify the change",
            "6. Inventory changed files",
            "7. Review the change",
            "8. Summarize the result",
        ],
        top=1.8,
        height=4.0,
    )
    add_footer(slide, "9")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Key Messages To Say Out Loud")
    add_bullets(
        slide,
        [
            "We are not asking Cursor to do everything. We are giving it a small task with constraints.",
            "Rules should capture non-obvious repo guidance, not generic engineering advice.",
            "We plan first when we want to control scope.",
            "We still verify with tests and build output.",
            "We still review the changed files before trusting the result.",
            "The goal is not more code. The goal is faster understanding and safer iteration.",
        ],
    )
    add_footer(slide, "10")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Fallback Shorter Exercise")
    add_bullets(
        slide,
        [
            "If time runs short, use a smaller feature: add `user.id` from `userIdentity.principalId`.",
            "Make it required in the schema and update `tests/test_mapping.py`.",
            "Skip the frontend UI change.",
            "This still demonstrates planning, implementation, verification, and review with less typing.",
        ],
        top=1.8,
    )
    add_footer(slide, "11")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Close")
    add_bullets(
        slide,
        [
            "Start with understanding, not generation.",
            "Keep tasks small and specific.",
            "Ask what changed, run checks, and review the result.",
            "Use AI where it helps the team move faster without lowering standards.",
        ],
        top=1.9,
    )
    add_footer(slide, "12")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(OUTPUT)
        print(f"Saved deck to {OUTPUT}")
    except PermissionError:
        prs.save(FALLBACK_OUTPUT)
        print(f"Saved deck to {FALLBACK_OUTPUT}")


if __name__ == "__main__":
    build_deck()
